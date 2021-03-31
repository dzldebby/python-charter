from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
import win32com.client
import pandas as pd
import os

# Fill in number of charts / slides you want to chart 
num_of_charts = 3

# Create presentation with 1 slide
prs = Presentation('chart-01.pptx')
# Select first slide 

data = pd.read_csv('data.csv', index_col=None)

df = {}
title = {}

# get data

col = []

for i in range(0,num_of_charts):
    df[str(i)] = data[i*7:4+(7)*i]
    df[str(i)].columns = range(df[str(i)].shape[1])
    df[str(i)].reset_index(drop=True, inplace=True)
    df[str(i)].columns = [df[str(i)].iloc[0][x] for x in range(0, len(data.columns))]
    # df[str(i)].columns = [df[str(i)].iloc[0][0], df[str(i)].iloc[0][1], df[str(i)].iloc[0][2], df[str(i)].iloc[0][3]]
    df[str(i)] = df[str(i)].drop(0, axis = 0)

# get first title

title["0"] = data.columns[0]


# get remainding titles 

for i in range(1,num_of_charts):
    title[str(i)] = data.iloc[i*6 + i -1][0]


# TABLE ----------------------------------

def create_table(data):
    # Create table objec 

    # 2nd placeholder     
    table_placeholder = slide.shapes[0]
    shape = table_placeholder.insert_table(rows=data.shape[0]+1, cols=data.shape[1])
    table = shape.table

    # Populate table header column text 

    for col in range(0, data.shape[1]):
        table.cell(0,col).text = str(data.columns[col])

    # Populate table data (except first column with the scale)
    for col in range(1,data.shape[1]):
        for row in range(0,data.shape[0]):
            table.cell(row+1,col).text = str(data.iloc[row][col]) + "%"
            font = table.cell(row+1,col).text_frame.paragraphs[0].font
            font.size = Pt(14)

    # Populate first column (scale data)
    for row in range(1,data.shape[0]+1):
        table.cell(row,0).text = str(data.iloc[row-1][0])
        font = table.cell(row,0).text_frame.paragraphs[0].font
        font.size = Pt(14)


# STACKED BAR CHART --------------------

def create_stackedbar(data):
    # Define chart data
    chart_data = CategoryChartData()

    # Transpose data, because that's how the stacked bar chart data requirement for ppt is 
    dataT = data.set_index('Scale').transpose()
    # Reverse order of columns, similar to 'Values in reverse order' in ppt 
    dataT = dataT[dataT.columns[::-1]]


    # For each column in the transposed data, add a series to the data
    for col_id, col in enumerate(dataT.columns):
        chart_data.add_series(dataT.columns[col_id], (dataT.iloc[:, col_id]))


    data_cat_list = []

    # Using the original data, create a list containing all the categories 
    for i in range(1,data.shape[1]):
        data_cat_list.append(data.columns[i])

    chart_data.categories = data_cat_list


    # Add chart to slide
    x, y, cx, cy = Cm(2), Cm(5), Cm(10), Cm(8)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
    )

    chart = graphic_frame.chart


    # data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '0"%"'

    # value_axis
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 100.0
    value_axis.has_minor_gridlines = False

    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(14)


def save_ppt():
    prs.save('chart-01.pptx')

def open_ppt():
    # Open ppt
    PptApp = win32com.client.Dispatch("Powerpoint.Application")
    PptApp.Visible = True
    PptApp.Presentations.Open(r'C:\Users\Debby\python-chart\chart-01.pptx')



def add_title(title):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        else:
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = title



# slide = prs.slides[0]

for i in range(0,num_of_charts):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide = prs.slides[i]
    create_table(df[str(i)])
    create_stackedbar(df[str(i)])
    add_title(title[str(i)])
    save_ppt()

open_ppt()
print("Charted!")

# create_table(df["0"])
# create_stackedbar(df["1"])
# add_title(title["1"])
# save_ppt()
# open_ppt()




